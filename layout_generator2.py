from pydantic import BaseModel, Field
from typing import List, Dict, Optional, Tuple
import json
from pptx import Presentation
import os
from slide_content_generator import generate_slide_content, Slide

class SlideContent(BaseModel):
    bullets: List[str]
    image_paths: List[str] = Field(default_factory=list)

class SlideContentForLayout(BaseModel):
    slide: int
    slide_title: str
    slide_content: SlideContent

class Layout(BaseModel):
    slide_dimensions: dict
    title_box: dict
    subtitle_box: dict
    bulleted: bool
    content_font_size: int
    text_box: dict
    image_paths: List[str]
    image_boxes: List[dict]
    title: str
    subtitle: str
    text: List[str]
    template_layout_index: int

layouts: List[Layout] = []

def get_template_layouts(template_path: str) -> List[Dict]:
    """Get available layouts from the PowerPoint template"""
    prs = Presentation(template_path)
    template_layouts = []
    
    for idx, layout in enumerate(prs.slide_layouts):
        layout_info = {
            "index": idx,
            "name": layout.name if hasattr(layout, 'name') else f"Layout {idx}",
            "placeholders": []
        }
        
        # Analyze placeholders in the layout
        for shape in layout.placeholders:
            placeholder_info = {
                "type": shape.placeholder_format.type,
                "idx": shape.placeholder_format.idx
            }
            layout_info["placeholders"].append(placeholder_info)
        
        template_layouts.append(layout_info)
    
    return template_layouts

def generate_layout_for_slide(slide_data: dict, template_layouts: List[Dict]) -> dict:
    """Generate layout for a single slide's data based on template"""
    # Get the template layout being used
    template_layout = next(
        (layout for layout in template_layouts if layout["index"] == slide_data["layout_id"]),
        template_layouts[1]  # Default to title and content layout
    )
    
    # Default slide dimensions (16:9 aspect ratio)
    slide_dimensions = {
        "width": 10,
        "height": 5.625
    }
    
    # Default title box
    title_box = {
        "x": 0.5,
        "y": 0.5,
        "width": 9,
        "height": 0.5,
        "font_size": 24,
        "padding": 0.1
    }
    
    # Default subtitle box
    subtitle_box = {
        "x": 0.5,
        "y": 1.2,
        "width": 9,
        "height": 0.5,
        "font_size": 18,
        "padding": 0.1
    }
    
    # Default text box
    text_box = {
        "layout": "left",
        "x": 0.5,
        "y": 1.8,
        "width": 4.5,
        "height": 3.5,
        "padding": 0.2
    }
    
    # Create image box if there are images
    image_boxes = []
    if slide_data["content"].get("image_path"):
        image_boxes.append({
            "layout": "right",
            "x": 5.5,
            "y": 1.8,
            "width": 4.5,
            "height": 3.5,
            "padding": 0.2
        })
    
    # Get image path
    image_paths = [slide_data["content"]["image_path"]] if slide_data["content"].get("image_path") else []
    
    return {
        "slide_dimensions": slide_dimensions,
        "title_box": title_box,
        "subtitle_box": subtitle_box,
        "bulleted": True,
        "content_font_size": 16,
        "text_box": text_box,
        "image_paths": image_paths,
        "image_boxes": image_boxes,
        "title": slide_data["title"],
        "subtitle": f"Slide {slide_data['slide_number']}",
        "text": slide_data["content"]["bullets"],
        "template_layout_index": slide_data["layout_id"]
    }

def get_slide_layout(slide_content: List[SlideContentForLayout], template_path: str):
    """Process all slides and generate layouts based on template"""
    global layouts
    layouts = []  # Reset layouts list
    
    # Get available layouts from template
    template_layouts = get_template_layouts(template_path)
    print(f"Found {len(template_layouts)} layouts in template")
    
    try:
        # Read slide layouts from JSON file
        with open('slide_layouts.json', 'r') as f:
            slide_layouts = json.load(f)
        print(f"Loaded {len(slide_layouts)} slide layouts from slide_layouts.json")
        
        # Generate layouts for each slide
        for slide_data in slide_layouts:
            layout_data = generate_layout_for_slide(slide_data, template_layouts)
            validated_layout = Layout(**layout_data)
            layouts.append(validated_layout)

        # Save layouts to file
        with open('generated_layouts.json', 'w') as f:
            json.dump([layout.dict() for layout in layouts], f, indent=2)
        print(f"✅ Generated layouts for {len(layouts)} slides and saved to generated_layouts.json")
        
    except FileNotFoundError:
        print("Error: slide_layouts.json not found")
        exit(1)
    except json.JSONDecodeError:
        print("Error: Invalid JSON in slide_layouts.json")
        exit(1)
    except Exception as e:
        print(f"Error: An unexpected error occurred: {str(e)}")
        exit(1)

# if __name__ == "__main__":
#     # Example usage
#     template_path = "available_templates/A.pptx"  # Update with your template path
#     get_slide_layout([], template_path)  # Empty list since we're using slide_layouts.json