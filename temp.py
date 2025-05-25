from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.util import Inches
import os
from openai import OpenAI
import json

client = OpenAI()

def extract_layout_info(prs):
    layout_specs = []
    for i, layout in enumerate(prs.slide_layouts):
        placeholders = []
        for shape in layout.placeholders:
            phf = shape.placeholder_format
            ph_type = PP_PLACEHOLDER(phf.type).name
            placeholders.append(f"{ph_type} {phf.idx}")
        layout_specs.append({
            "layout_idx": i,
            "layout_name": layout.name,
            "placeholders": placeholders
        })
    return layout_specs

def call_llm_generate_slides(prompt, function_schema):
    response = client.chat.completions.create(
        model="gpt-3.5-turbo",
        messages=[
            {"role": "system", "content": "You are a slide generation assistant."},
            {"role": "user", "content": prompt}
        ],
        functions=[function_schema],
        function_call={"name": "generate_slide_content"}
    )
    return json.loads(response.choices[0].message.function_call.arguments)

def create_presentation(slide_data, prs):
    for slide_info in slide_data["slides"]:
        layout = prs.slide_layouts[slide_info["layout_idx"]]
        slide = prs.slides.add_slide(layout)
        if "title" in slide_info and slide.shapes.title:
            slide.shapes.title.text = slide_info["title"]
        if "content" in slide_info:
            for shape in slide.placeholders:
                if shape.placeholder_format.type == PP_PLACEHOLDER.BODY:
                    shape.text = slide_info["content"]
                    break
        if "image_path" in slide_info and os.path.exists(slide_info["image_path"]):
            left = Inches(1)
            top = Inches(2)
            slide.shapes.add_picture(slide_info["image_path"], left, top, height=Inches(3))
        if "speaker_notes" in slide_info:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = slide_info["speaker_notes"]

def generate_prompt(input_item, layout_specs):
    layout_text = json.dumps(layout_specs, indent=2)
    return f"""
    Given the detailed input, generate a list of slides with layout_idx, title, content, optional image_path, and speaker_notes.

    Available layouts:
    {layout_text}

    Input:
    {json.dumps(input_item, indent=2)}
    """

if __name__ == "__main__":
    prs = Presentation("available_templates/A.pptx")
    layout_specs = extract_layout_info(prs)

    with open("presentation_data.json", "r", encoding="utf-8") as f:
        input_data_list = json.load(f)

    prompt = generate_prompt(input_data_list[0], layout_specs)
    print(prompt)

    function_schema = {
        "name": "generate_slide_content",
        "description": "Generates slide content based on the provided input.",
        "parameters": {
            "type": "object",
            "properties": {
                "slides": {
                    "type": "array",
                    "items": {
                        "type": "object",
                        "properties": {
                            "layout_idx": {"type": "integer"},
                            "title": {"type": "string"},
                            "content": {"type": "string"},
                            "image_path": {"type": "string"},
                            "speaker_notes": {"type": "string"}
                        },
                        "required": ["layout_idx", "title", "content"]
                    }
                }
            },
            "required": ["slides"]
        }
    }

    # structured_output = call_llm_generate_slides(prompt, function_schema)
    # print(structured_output)

    # for idx, input_item in enumerate(input_data_list):
    #     prompt = generate_prompt(input_item, layout_specs)
    #     structured_output = call_llm_generate_slides(prompt, function_schema)
    #     create_presentation(structured_output, prs)

    # prs.save("generated_airbnb_slides.pptx")