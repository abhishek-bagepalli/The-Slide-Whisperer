from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import MSO_AUTO_SIZE
import os

def create_slide_from_content(template_path, output_path, slides_data):
    """
    Create a presentation with multiple slides using the provided content and template.
    
    Args:
        template_path (str): Path to the PowerPoint template (.pptx)
        output_path (str): Path where the new presentation will be saved
        slides_data (list): List of dictionaries containing slide content and layout information
    """
    # Load the template
    prs = Presentation(template_path)
    
    # Process each slide
    for slide_data in slides_data:
        # Get the layout based on layout_id
        layout_id = slide_data['layout_id']
        slide_layout = prs.slide_layouts[layout_id]
        slide = prs.slides.add_slide(slide_layout)
        
        # First, find all placeholders and their types
        placeholders = {}
        for shape in slide.placeholders:
            placeholder_type = shape.placeholder_format.type
            placeholder_idx = shape.placeholder_format.idx
            placeholders[placeholder_idx] = shape
        
        # Process each content item based on its mapping
        for item in slide_data['mapping']:
            placeholder_type = item['placeholder_type']
            placeholder_index = item['placeholder_index']
            
            if placeholder_index in placeholders:
                shape = placeholders[placeholder_index]
                
                # Handle title
                if item['content_type'] == 'title':
                    shape.text = item['value']
                    text_frame = shape.text_frame
                    text_frame.word_wrap = True
                    text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
                    # Set font size for title - larger for title slide
                    for paragraph in text_frame.paragraphs:
                        if layout_id == 0:  # Title slide
                            paragraph.font.size = Pt(44)  # Larger font for title slide
                        else:
                            paragraph.font.size = Pt(20)  # Normal font for other slides
                
                # Handle subtitle (for title slide)
                elif item['content_type'] == 'subtitle' and layout_id == 0:
                    shape.text = item['value']
                    text_frame = shape.text_frame
                    text_frame.word_wrap = True
                    text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
                    # Set font size for subtitle
                    for paragraph in text_frame.paragraphs:
                        paragraph.font.size = Pt(32)  # Larger font for subtitle
                
                # Handle bullets
                elif item['content_type'] == 'bullets':
                    text_frame = shape.text_frame
                    text_frame.clear()
                    text_frame.word_wrap = True
                    text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
                    
                    # Handle bullet points - now accepting both string and list inputs
                    bullet_points = item['value']
                    if isinstance(bullet_points, str):
                        bullet_points = [point.strip() for point in bullet_points.split('\n') if point.strip()]
                    
                    for i, bullet in enumerate(bullet_points):
                        if i == 0:
                            p = text_frame.paragraphs[0]
                        else:
                            p = text_frame.add_paragraph()
                        p.text = bullet
                        p.level = 0
                        # Set font size for bullets
                        p.font.size = Pt(20)
                
                # Handle image
                elif item['content_type'] == 'image_path':
                    image_path = './images/' + item['value']
                    if os.path.exists(image_path):
                        try:
                            left = shape.left
                            top = shape.top
                            width = shape.width
                            height = shape.height
                            slide.shapes.add_picture(image_path, left, top, width, height)
                        except Exception as e:
                            print(f"⚠️ Error adding image: {str(e)}")
                    else:
                        print(f"⚠️ Image not found: {image_path}")
                
                # Handle speaker notes
                elif item['content_type'] == 'speaker_notes':
                    notes_slide = slide.notes_slide
                    notes_slide.notes_text_frame.text = item['value']
                    # Set font size for speaker notes
                    for paragraph in notes_slide.notes_text_frame.paragraphs:
                        paragraph.font.size = Pt(20)
    
    # Save the presentation
    prs.save(output_path)
    print(f"✅ Presentation created with {len(slides_data)} slides and saved to: {output_path}")

if __name__ == "__main__":
    # Example usage
    template_path = "available_templates/A.pptx"
    output_path = "generated_presentation.pptx"
    
    slides_data = [
        {
            "slide_number": 1,
            "layout_id": 0,  # Title slide layout
            "layout_name": "Title Slide",
            "mapping": [
                {
                    "content_type": "title",
                    "value": "Airbnb: A Journey of Innovation and Growth",
                    "placeholder_type": "TITLE",
                    "placeholder_index": 0
                },
                {
                    "content_type": "subtitle",
                    "value": "Exploring the Evolution of the World's Leading Lodging Marketplace",
                    "placeholder_type": "SUBTITLE",
                    "placeholder_index": 1
                }
            ]
        },
        {
            "slide_number": 2,
            "layout_id": 1,
            "layout_name": "Title and content with image right",
            "mapping": [
                {
                    "content_type": "title",
                    "value": "Investment Milestone",
                    "placeholder_type": "TITLE",
                    "placeholder_index": 0
                },
                {
                    "content_type": "bullets",
                    "value": [
                        "Securing $7.2 million Series A funding from Greylock Partners and Sequoia Capital",
                        "Facing competition from Wimdu and strategic decisions to maintain competitive edge",
                        "Balancing growth and mission-driven focus"
                    ],
                    "placeholder_type": "OBJECT",
                    "placeholder_index": 1
                },
                {
                    "content_type": "image_path",
                    "value": "reid_hoffman's_investment_reflections.jpg",
                    "placeholder_type": "PICTURE",
                    "placeholder_index": 13
                },
                {
                    "content_type": "speaker_notes",
                    "value": "Investment challenges and competitive dynamics shaped Airbnb's strategic direction.",
                    "placeholder_type": "SLIDE_NUMBER",
                    "placeholder_index": 11
                }
            ]
        }
    ]
    
    create_slide_from_content(template_path, output_path, slides_data)