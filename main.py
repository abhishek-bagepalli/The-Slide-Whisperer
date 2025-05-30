import json
from document_parser import extract_text_and_tables
from tools import clear_images_folder, PresentationSummarizer
from text_chunker import chunk_text
from multi_document_rag import MultiDocumentRAG
from multimodal_rag import build_image_index
from get_image_from_web import search_and_download_image_from_web
from tools import get_best_image
from slide_content_generator import generate_slide_content
from tools import update_image_dimensions
from slide_content_generator import get_llm_friendly_layouts
from presentation_pipeline import run_presentation_pipeline
from pptx import Presentation

def main():
    # === PARAMETERS ===
    document_path = "docs/Airbnb Case.pdf"  # Change as needed
    image_folder = "./images"
    min_chunk_size = 1000
    max_chunk_size = 5000
    minimum_slides = 7
    chosen_template = "available_templates/A.pptx"
    output_path = "outputs/slide_whisper3.pptx"

    print("🚀 Starting presentation generation process...")

    # 1. Clear images folder
    print("🗑️ Clearing existing images folder...")
    clear_images_folder()

    # 2. Parse document and extract images/metadata
    print(f"📄 Processing document: {document_path}")
    print("🔎 Initializing RAG and answering document queries...")
    rag = MultiDocumentRAG()
    rag.process_documents([document_path])

    with open('document_parsed.json', 'r') as f:
        json_result = json.load(f)

    # 3. Extract text and tables
    print("📝 Extracting text and tables from parsed document...")
    text = extract_text_and_tables(json_result)

    # 4. Chunk the text
    print("✂️ Chunking text...")
    text_chunks = chunk_text(text, min_chunk_size=min_chunk_size, max_chunk_size=max_chunk_size)
    print(f"Generated {len(text_chunks)} text chunks.")

    # 5. Summarize each chunk and ideate on visualizations/queries
    print("🤖 Summarizing chunks and ideating visualizations...")
    summarizer = PresentationSummarizer()
    summarizer.clear_history()
    summaries = []
    for i, chunk in enumerate(text_chunks):
        previous_chunk = text_chunks[i-1] if i > 0 else None
        summary = summarizer.summarize_text(chunk, previous_chunk)
        summaries.append(summary)
        print(f"Chunk {i+1}/{len(text_chunks)} summarized.")

    # 6. Use MultiDocumentRAG to answer document queries
    query_results_from_document = []
    for summary in summaries:
        doc_queries = summary.additional_information_needed["document_queries"]
        
        for query in doc_queries:
            print(f"\nProcessing query: {query}")
            query_result = {"query": query,"document_response": None}
            
            responses = rag.get_exact_content(query, 1)
            
            for response in responses:
                print("Found relevant content in document")
                query_result["document_response"] = response['content']
            
            query_results_from_document.append(query_result)

    # 7. Build image index and retrieve/download images for visualizations
    print("🖼️ Building image index and retrieving images...")
    image_index = build_image_index(image_folder)

    # 8. Create presentation data
    presentation_data = []
    for summary in summaries:
        summary_data = {
            "detailed_summary": summary.detailed_summary,
            "key_visualizations": {
                "charts": summary.key_visualizations["charts"],
                "images": summary.key_visualizations["images"],
                "retrived_image_paths_charts": [],
                "retrived_image_paths_images": []
            },
            "retrieved_content_from_document": []
        }
        
        # Add document responses
        for query in summary.additional_information_needed["document_queries"]:
            for result in query_results_from_document:
                if result["query"] == query:
                    summary_data["retrieved_content_from_document"].append({
                        "query": query,
                        "response": result["document_response"]
                    })
        
        # Process charts
        for query in summary.key_visualizations['charts']:
            image_path, confidence = get_best_image(query, image_index)
            if confidence > 0.3:
                summary_data["key_visualizations"]["retrived_image_paths_charts"].append(image_path)
            else:
                image_path = search_and_download_image_from_web(query)
                if image_path:
                    summary_data["key_visualizations"]["retrived_image_paths_charts"].append(image_path)
        
        # Process images
        for query in summary.key_visualizations['images']:
            image_path, confidence = get_best_image(query, image_index)
            if confidence > 0.3:
                summary_data["key_visualizations"]["retrived_image_paths_images"].append(image_path)
            else:
                image_path = search_and_download_image_from_web(query)
                if image_path:
                    summary_data["key_visualizations"]["retrived_image_paths_images"].append(image_path)
        
        presentation_data.append(summary_data)

    # Save to JSON
    with open('presentation_data.json', 'w') as f:
        json.dump(presentation_data, f, indent=2)
    print("✅ Presentation data saved to presentation_data.json")

    with open('presentation_data.json', 'r') as f:
        presentation_data = json.load(f)

    slides, metadata = generate_slide_content(presentation_data, minimum_slides)

    # Save metadata to a JSON file
    metadata_dict = metadata.model_dump()
    with open('metadata.json', 'w') as f:
        json.dump(metadata_dict, f, indent=2)
    print("✅ Metadata saved to metadata.json")

    # Read metadata from JSON file
    with open('metadata.json', 'r') as f:
        metadata_dict = json.load(f)
    

    print("Generated slides:", slides)
    print("Generated metadata:", metadata)

    # Convert Slide objects to dictionaries before JSON serialization
    slide_content_dict = [slide.model_dump() for slide in slides]
    with open('slide_content.json', 'w') as f:
        json.dump(slide_content_dict, f, indent=2)

    with open('slide_content.json', 'r') as f:
        slide_content = json.load(f)

    updated_slide_content = update_image_dimensions(slide_content)

    layout_specs = get_llm_friendly_layouts(chosen_template)

    print(layout_specs)

    run_presentation_pipeline(chosen_template, output_path, 'slide_content.json', layout_specs)

    # Add title and subtitle to first slide
    prs = Presentation(output_path)
    first_slide = prs.slides[0]
    
    # Clear existing content
    for shape in first_slide.shapes:
        if shape.has_text_frame:
            shape.text_frame.clear()
    
    # Add title
    title_shape = first_slide.shapes.title
    if title_shape:
        title_shape.text = metadata_dict['title']
    
    # Add subtitle
    subtitle_shape = first_slide.placeholders[1]  # Usually the subtitle placeholder
    if subtitle_shape:
        subtitle_shape.text = metadata_dict['subtitle']
    
    # Save the modified presentation
    prs.save(output_path)
    print("✅ Added title and subtitle to first slide")

if __name__ == "__main__":
    main()
