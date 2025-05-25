# slide_tools.py

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os
import json

# Create or load a presentation
PPTX_PATH = "airbnb_v1.pptx"
# PPTX_PATH_CLEAN = "loreal_presentation_v1_clean.pptx"

prs = Presentation(PPTX_PATH)

SLIDE_WIDTH = prs.slide_width.inches
SLIDE_HEIGHT = prs.slide_height.inches

def save_presentation(path):
    prs.save(path)

def is_slide_empty(slide):
    # Check if all placeholders are empty
    for shape in slide.shapes:
        if shape.has_text_frame and shape.text_frame.text.strip():
            return False
        if shape.shape_type == MSO_SHAPE.PICTURE:
            return False
    return True

def delete_empty_slides(prs):
    # Create a list of slides to delete
    slides_to_delete = []
    for i, slide in enumerate(prs.slides):
        if is_slide_empty(slide):
            slides_to_delete.append(i)
    
    # Delete slides in reverse order to avoid index shifting
    for i in sorted(slides_to_delete, reverse=True):
        xml_slides = prs.slides._sldIdLst
        xml_slides.remove(xml_slides[i])
        print(f"Deleted slide at index {i} (empty)")

def create_slide_from_template(spec_arg, i):
    # prs.slide_height = Inches(spec_arg['slide_dimensions']['width'])
    # prs.slide_height = Inches(spec_arg['slide_dimensions']['height'])

    title = spec_arg.get("title")
    subtitle = spec_arg.get("subtitle")
    text_items = spec_arg.get("text", [])
    is_bulleted = spec_arg.get("bulleted", False)
    font_size = spec_arg.get("content_font_size", 24)

    text_box = spec_arg.get("text_box", {})
    image_paths = spec_arg.get("image_paths")
    image_boxes = spec_arg.get("image_boxes", {})

    padding = 0.2  # default padding
    text_padding = text_box.get("padding", padding)
    # image_padding = image_boxes[0].get("padding", padding)

    title_box = spec_arg.get("title_box", {})
    subtitle_box = spec_arg.get("subtitle_box", {})
    title_padding = 0.2
    layout_index = spec_arg.get("slide_layout", 3)

    # If this is the first slide, modify the existing title slide
    if i == 0:
        print('should be hit only once')
        slide = prs.slides[0]  # Access the existing first slide
        
        # Set title
        for placeholder in slide.placeholders:
            if placeholder.placeholder_format.type == PP_PLACEHOLDER.CENTER_TITLE:
                placeholder.text = title
                break

        # Set subtitle if it exists
        if subtitle:
            for placeholder in slide.placeholders:
                if placeholder.placeholder_format.type == PP_PLACEHOLDER.SUBTITLE:
                    placeholder.text = subtitle
                    break
    else:
        # For subsequent slides, create new slides
        slide_layout = prs.slide_layouts[layout_index]
        slide = prs.slides.add_slide(slide_layout)

        for placeholder in slide.placeholders:
            if placeholder.placeholder_format.type == PP_PLACEHOLDER.TITLE:
                placeholder.text = title
                break

        valid_content_types = {
            PP_PLACEHOLDER.BODY,
            PP_PLACEHOLDER.OBJECT,
            PP_PLACEHOLDER.PICTURE
        }

        content_placeholders = [
            ph for ph in slide.placeholders
            if ph.placeholder_format.type in valid_content_types
        ]

        # Insert bullet points into the first content placeholder (if available)
        if len(text_items) > 0 and len(content_placeholders) >= 1:
            content_ph = content_placeholders[0]
            content_ph.text = ""  # Clear default text
            if is_bulleted:
                for point in text_items:
                    p = content_ph.text_frame.add_paragraph()
                    p.text = point
                    p.level = 0
                    p.font.size = Pt(font_size)
            else:
                content_ph.text_frame.text = "\n".join(text_items)
                content_ph.text_frame.paragraphs[0].font.size = Pt(font_size)

        # Insert image into second placeholder (if available)
        if len(image_paths) > 0 and len(content_placeholders) >= 2:
            image_ph = content_placeholders[1]
            try:
                print('hit 3')
                slide.shapes.add_picture(image_paths[0], image_ph.left, image_ph.top, width=image_ph.width)
            except Exception as e:
                print(f"⚠️ Could not insert image: {e}")

    save_presentation(PPTX_PATH)
    return {
        "status": "success",
        "message": "Slide created with layout positioning and formatting.",
        "title": title,
        "slide_number": len(prs.slides)
    }


with open("generated_layouts.json", "r") as f:
    slide_specs = json.load(f)

print(slide_specs)

i = 0
for spec in slide_specs:
    result = create_slide_from_template(spec, i)
    i += 1
    print(f"✅ Slide created: {result}")
