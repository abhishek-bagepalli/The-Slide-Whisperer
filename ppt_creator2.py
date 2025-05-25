from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import MSO_AUTO_SIZE
import json
import os

# Mapping of placeholder type numbers to their names
PLACEHOLDER_TYPES = {
    1: "TITLE",
    2: "SUBTITLE",
    3: "BODY",
    4: "TITLE_AND_BODY",
    5: "TITLE_AND_SUBTITLE",
    6: "TITLE_AND_VERTICAL_BODY",
    7: "OBJECT",
    8: "VERTICAL_TITLE_AND_TEXT",
    9: "VERTICAL_TITLE_AND_VERTICAL_TEXT",
    10: "TITLE_ONLY",
    11: "SECTION_HEADER",
    12: "TWO_CONTENT",
    13: "PICTURE",
    14: "TWO_OBJECTS",
    15: "TWO_OBJECTS_AND_OBJECT",
    16: "TWO_OBJECTS_AND_TEXT",
    17: "TWO_OBJECTS_OVER_TEXT",
    18: "PICTURE",
    19: "TEXT_AND_CLIP_ART",
    20: "CLIP_ART_AND_TEXT",
    21: "TITLE_AND_CLIP_ART",
    22: "CLIP_ART_AND_VERTICAL_TEXT",
    23: "VERTICAL_TITLE_AND_TEXT_OVER_CHART",
    24: "TEXT_AND_CHART",
    25: "CHART_AND_TEXT",
    26: "TITLE_AND_CHART",
    27: "CHART_AND_TEXT",
    28: "TEXT_AND_CLIP_ART",
    29: "CLIP_ART_AND_TEXT",
    30: "TITLE_AND_CLIP_ART",
    31: "CLIP_ART_AND_VERTICAL_TEXT",
    32: "VERTICAL_TITLE_AND_TEXT_OVER_CHART",
    33: "TEXT_AND_CHART",
    34: "CHART_AND_TEXT",
    35: "TITLE_AND_CHART",
    36: "CHART_AND_TEXT"
}

def get_available_layouts(prs: Presentation) -> list:
    """Get a list of available layout names from the presentation template."""
    layouts = []
    for layout in prs.slide_layouts:
        name = layout.name if hasattr(layout, 'name') else f"Layout {len(layouts)}"
        layouts.append(name)
    return layouts

def create_presentation_from_template(template_path, output_path, layouts_json_path):
    """
    Create a presentation using an existing template and layout data.
    
    Args:
        template_path (str): Path to the PowerPoint template (.pptx)
        output_path (str): Path where the new presentation will be saved
        layouts_json_path (str): Path to the JSON file containing slide layouts
    """
    # Image name mapping
    image_name_mapping = {
        "airbnb_early_timeline_infographic.jpg": "airbnb's_early_days_setup_with_airbeds.jpg",
        "obama_o's_and_cap'n_mccain's_cereal_boxes.jpg": "designs_of_obama_o's_and_cap'n_mccain's_cereal_box.jpg",
        "airbnb_rebranding_process.jpg": "airbnb's_growth_in_number_of_hosts_and_listings_ov.jpg",
        "airbnb_hosts_interacting_with_photographers_for_property_photos.jpg": "high-quality_property_photos_before_and_after_airb.jpg",
        "reid_hoffman_presenting_airbnb_investment_to_greyl.jpg": "reid_hoffman_reflecting_on_airbnb_investment_at_gr.jpg",
        "airbnb_local_offices_setup_in_germany.jpg": "local_office_opening_ceremony_in_hamburg.jpg",
        "airbnb_professional_photography_session.jpg": "professional_photography_service_launch_event_at_a.jpg",
        "sequoia_capital_investment_in_airbnb.jpg": "airbnb_funding_rounds_and_investments_over_time.jpg",
        "comparison_of_funding_rounds_in_airbnb's_growth_jo.jpg": "comparison_of_funding_rounds_between_airbnb_and_wi.jpg"
    }
    
    # Load the template
    prs = Presentation(template_path)
    
    # Get available layouts
    available_layouts = get_available_layouts(prs)
    print("\n=== Available Layouts in Template ===")
    for i, layout in enumerate(available_layouts):
        print(f"{i}: {layout}")
    
    # Load the layouts
    with open(layouts_json_path, 'r') as f:
        layouts = json.load(f)
    
    # Process each layout and create slides
    for layout_data in layouts:
        print(f"\n=== Processing Slide: {layout_data['title']} ===")
        print(f"Template Layout Index: {layout_data.get('template_layout_index', 1)}")
        
        # Get the template layout index
        template_layout_index = layout_data.get('template_layout_index', 1)
        slide_layout = prs.slide_layouts[template_layout_index]
        slide = prs.slides.add_slide(slide_layout)
        
        # First, find all placeholders and their types
        placeholders = {}
        print("\nFound Placeholders:")
        for shape in slide.placeholders:
            placeholder_type = shape.placeholder_format.type
            placeholder_idx = shape.placeholder_format.idx
            placeholders[placeholder_type] = shape
            type_name = PLACEHOLDER_TYPES.get(placeholder_type, f"UNKNOWN ({placeholder_type})")
            print(f"- Type {type_name}, Index {placeholder_idx}")
        
        # Handle title
        if 1 in placeholders:  # Title
            title_shape = placeholders[1]
            print(f"\nAdding Title:")
            print(f"- Content: {layout_data['title']}")
            title_shape.text = layout_data['title']
            
            # Enable text wrapping and auto-sizing for title
            text_frame = title_shape.text_frame
            text_frame.word_wrap = True
            text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
            
            if 'font_size' in layout_data['title_box']:
                for paragraph in text_frame.paragraphs:
                    paragraph.font.size = Pt(layout_data['title_box']['font_size'])
                print(f"- Font Size: {layout_data['title_box']['font_size']}")
            print("✅ Title added successfully")
        else:
            print("⚠️ No title placeholder found")
        
        # Handle subtitle
        if 2 in placeholders:  # Subtitle
            subtitle_shape = placeholders[2]
            print(f"\nAdding Subtitle:")
            print(f"- Content: {layout_data['subtitle']}")
            subtitle_shape.text = layout_data['subtitle']
            
            # Enable text wrapping and auto-sizing for subtitle
            text_frame = subtitle_shape.text_frame
            text_frame.word_wrap = True
            text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
            
            if 'font_size' in layout_data['subtitle_box']:
                for paragraph in text_frame.paragraphs:
                    paragraph.font.size = Pt(layout_data['subtitle_box']['font_size'])
                print(f"- Font Size: {layout_data['subtitle_box']['font_size']}")
            print("✅ Subtitle added successfully")
        else:
            print("⚠️ No subtitle placeholder found")
        
        # Handle body content - try both BODY (3) and OBJECT (7) placeholders
        content_shape = None
        if 3 in placeholders:  # Body
            content_shape = placeholders[3]
            print("\nUsing BODY placeholder for content")
        elif 7 in placeholders:  # Object
            content_shape = placeholders[7]
            print("\nUsing OBJECT placeholder for content")
        
        if content_shape:
            print(f"\nAdding Body Content:")
            print(f"- Number of bullet points: {len(layout_data['text'])}")
            text_frame = content_shape.text_frame
            text_frame.clear()
            
            # Enable text wrapping and auto-sizing for body content
            text_frame.word_wrap = True
            text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
            
            # Add bullets
            for i, bullet in enumerate(layout_data['text']):
                if i == 0:
                    # Use the first paragraph that's created by clear()
                    p = text_frame.paragraphs[0]
                else:
                    # Add new paragraphs for subsequent bullets
                    p = text_frame.add_paragraph()
                
                p.text = bullet
                p.level = 0
                print(f"- Bullet {i+1}: {bullet[:50]}...")
                
                if 'content_font_size' in layout_data:
                    p.font.size = Pt(layout_data['content_font_size'])
                    print(f"- Font Size: {layout_data['content_font_size']}")
            print("✅ Body content added successfully")
        else:
            print("⚠️ No content placeholder found (neither BODY nor OBJECT)")
        
        # Handle picture - check both type 13 and type 18
        picture_types = [13, 18]  # Support both common picture placeholder types
        for pic_type in picture_types:
            if pic_type in placeholders and layout_data['image_paths']:
                picture_shape = placeholders[pic_type]
                img_path = layout_data['image_paths'][0]  # Take first image
                
                print(f"\nAdding Image:")
                print(f"- Image path: {img_path}")
                
                # Try different possible image paths
                possible_paths = [
                    img_path,  # Try direct path first
                    os.path.join("images", img_path),  # Try in images folder
                    os.path.join("images", img_path.lower()),  # Try lowercase version
                    os.path.join("images", img_path.replace(" ", "_")),  # Try with underscores
                    os.path.join("images", img_path.replace(" ", "_").lower()),  # Try both
                    os.path.join("images", image_name_mapping.get(img_path, img_path))  # Try mapped name
                ]
                
                image_found = False
                for path in possible_paths:
                    if os.path.exists(path):
                        try:
                            # Get placeholder dimensions
                            left = picture_shape.left
                            top = picture_shape.top
                            width = picture_shape.width
                            height = picture_shape.height
                            
                            # Add the image
                            slide.shapes.add_picture(path, left, top, width, height)
                            print(f"✅ Image added successfully from: {path}")
                            image_found = True
                            break
                        except Exception as e:
                            print(f"⚠️ Error adding image from {path}: {str(e)}")
                
                if not image_found:
                    print(f"⚠️ Image not found in any of these locations:")
                    for path in possible_paths:
                        print(f"  - {path}")
                break  # Stop after trying to add image to first available picture placeholder
    
    # Save the presentation
    prs.save(output_path)
    print(f"\n✅ Presentation saved to: {output_path}")

if __name__ == "__main__":
    template_path = "available_templates/A.pptx"
    output_path = "generated_presentation.pptx"
    layouts_json_path = "generated_layouts.json"
    
    create_presentation_from_template(template_path, output_path, layouts_json_path)
