from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.shapes import PP_PLACEHOLDER
import json
import os

def create_presentation_from_layouts(template_path, output_path, layouts_json_path):
    """
    Create a presentation using an existing template and layout data from slide_layouts.json.
    
    Args:
        template_path (str): Path to the PowerPoint template (.pptx)
        output_path (str): Path where the new presentation will be saved
        layouts_json_path (str): Path to the JSON file containing slide layouts
    """
    # Load the template
    prs = Presentation(template_path)
    
    # Load the layouts
    with open(layouts_json_path, 'r') as f:
        layouts = json.load(f)
    
    # Process each layout and create slides
    for layout_data in layouts:
        print(f"\n=== Processing Slide: {layout_data['title']} ===")
        
        # Get the layout based on layout_id
        layout_id = layout_data['layout_id']
        slide_layout = prs.slide_layouts[layout_id]
        slide = prs.slides.add_slide(slide_layout)
        
        # First, find all placeholders and their types
        placeholders = {}
        for shape in slide.placeholders:
            placeholder_type = shape.placeholder_format.type
            placeholder_idx = shape.placeholder_format.idx
            placeholders[placeholder_type] = shape
        
        # Handle title
        if 1 in placeholders:  # Title
            title_shape = placeholders[1]
            title_shape.text = layout_data['title']
        
        # Handle content (bullets)
        content_shape = None
        for ph_type in [2, 3]:  # BODY and OBJECT placeholders
            if ph_type in placeholders:
                content_shape = placeholders[ph_type]
                break
        
        if content_shape and 'content' in layout_data and 'bullets' in layout_data['content']:
            text_frame = content_shape.text_frame
            text_frame.clear()
            
            # Add bullets
            for bullet in layout_data['content']['bullets']:
                p = text_frame.add_paragraph()
                p.text = bullet
                p.level = 0
        
        # Handle image
        if 'content' in layout_data and 'image_path' in layout_data['content']:
            image_path = layout_data['content']['image_path']
            
            # Check if image exists in current directory or images folder
            if not os.path.exists(image_path):
                image_path = os.path.join("images", image_path)
            
            if os.path.exists(image_path):
                try:
                    # Find picture placeholder
                    for ph_type in [13, 18]:  # PICTURE and MEDIA placeholders
                        if ph_type in placeholders:
                            picture_shape = placeholders[ph_type]
                            # Get placeholder dimensions
                            left = picture_shape.left
                            top = picture_shape.top
                            width = picture_shape.width
                            height = picture_shape.height
                            
                            # Add the image
                            slide.shapes.add_picture(image_path, left, top, width, height)
                            break
                except Exception as e:
                    print(f"⚠️ Error adding image: {str(e)}")
            else:
                print(f"⚠️ Image not found: {image_path}")
    
    # Save the presentation
    prs.save(output_path)
    print(f"\n✅ Presentation saved to: {output_path}")

if __name__ == "__main__":
    template_path = "available_templates/A.pptx"
    output_path = "generated_presentation.pptx"
    layouts_json_path = "slide_layouts.json"
    
    create_presentation_from_layouts(template_path, output_path, layouts_json_path) 