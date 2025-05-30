from typing import List, Optional, Tuple, Union
import json
from pydantic import BaseModel
from openai import OpenAI
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.shapes import PP_PLACEHOLDER

class ContentMapping(BaseModel):
    content_type: str
    value: Union[str, List[str]]
    placeholder_type: str
    placeholder_index: int

class SlideContent(BaseModel):
    bullets: List[str]
    speaker_notes: str
    image_paths: List[str] = []  # Default empty list to prevent validation error

class Slide(BaseModel):
    slide: int
    slide_title: str
    slide_content: SlideContent

class PresentationMetadata(BaseModel):
    title: str
    subtitle: str

class PresentationResponse(BaseModel):
    metadata: PresentationMetadata
    slides: List[Slide]

def generate_slide_content(presentation_data, minimum_slides=7) -> Tuple[List[Slide]]:
    client = OpenAI()
    minimum_slides = min(minimum_slides, 7)
    prompt = f"""
    Strict Instructions:

    minimum_slides: {minimum_slides}

    First, create a title and subtitle for the presentation.

    Provide a title and subtitle for the presentation.

    Then, think about all the `detailed_summary` and details in the `retrieved_content_from_document` and how they can be arranged to determine the overall story of the presentation.
    Then, use the information in the `detailed_summary` and the `response` in the `retrieved_content_from_document` to create informative and clear bullet points for each slide. Each bullet point should be detailed in itself.
    Ensure there are a maximum of 3 bullets per slide, if there is only one text box. If there are multiple text boxes, then there should be a maximum of 4 bullets per slide.
    Then, map all the key visualizations to appropriate slides using the provided image paths. 
    You are encouraged to use images. Ensure the image paths are taken from the `retrived_image_paths_charts` and `retrived_image_paths_images`. 
    Then, create a JSON object with the following structure:
    {{
        "metadata": {{
            "title": "Title of the presentation",
            "subtitle": "Subtitle of the presentation"
        }},

        "slides": [
            {{
                "slide": number,
                "slide_title": "Title of the slide",
                "slide_content": {{
                    "bullets": ["bullet point 1", "bullet point 2"], 
                    "speaker_notes": "Detailed speaker notes for this slide",
                    "image_paths": ["image_path_1","image_path_2"]
                }}
            }}
        ]
    }}
    {json.dumps(presentation_data, indent=2)}
"""
    response = client.chat.completions.create(
        model="gpt-3.5-turbo",
        messages=[
            {"role": "system", "content": "You are a presentation expert who creates detailed, informative section content in JSON format with specific examples and thorough analysis."},
            {"role": "user", "content": prompt}
        ],
        temperature=0.3,
        max_tokens=4000
    )
    try:
        content = response.choices[0].message.content.strip()
        content = content.replace('```json', '').replace('```', '').rstrip(',')
        raw_content = json.loads(content)
        
        # Validate metadata and slides using Pydantic
        slides = [Slide(**slide) for slide in raw_content["slides"]]
        metadata = PresentationMetadata(**raw_content["metadata"])
        return slides, metadata
        
    except json.JSONDecodeError as e:
        print(f"Error: Could not parse response as JSON. Error details: {str(e)}")
        print("Raw response content:", response.choices[0].message.content)
        return None, None
    except Exception as e:
        print(f"Error: Could not validate content. Error details: {str(e)}")
        return None, None

def get_available_layouts(template_path = "templates/A.pptx"):
    
    prs = Presentation(template_path)
    
    # Iterate through slide layouts
    for i, layout in enumerate(prs.slide_layouts):
        print(f"Layout {i}: {layout.name}")
        for shape in layout.placeholders:
            if shape.is_placeholder:
                phf = shape.placeholder_format
                print(f"  - Name: {shape.name}")
                print(f"    Type: {phf.type}")  # e.g., TITLE, BODY, PICTURE
                print(f"    Index (idx): {phf.idx}")
                print(f"    Placeholder Type: {PP_PLACEHOLDER(phf.type).name}")

def emu_to_inches(emu):
    return round(emu / 914400, 2)

def get_llm_friendly_layouts(pptx_path):
    prs = Presentation(pptx_path)
    layouts_info = []

    for layout_id, layout in enumerate(prs.slide_layouts):
        layout_data = {
            "layout_id": layout_id,
            "layout_name": layout.name,
            "placeholders": []
        }

        for shape in layout.placeholders:
            if shape.is_placeholder:
                phf = shape.placeholder_format
                try:
                    ph_type = PP_PLACEHOLDER(phf.type).name
                except ValueError:
                    ph_type = "UNKNOWN"

                layout_data["placeholders"].append({
                    "name": shape.name,
                    "placeholder_type": ph_type,
                    "index": phf.idx,
                    "position": {
                        "left_in": emu_to_inches(shape.left),
                        "top_in": emu_to_inches(shape.top)
                    },
                    "size": {
                        "width_in": emu_to_inches(shape.width),
                        "height_in": emu_to_inches(shape.height)
                    }
                })

        layouts_info.append(layout_data)

    return layouts_info
# print(get_llm_friendly_layouts("available_templates/A.pptx"))

def build_prompt_with_placeholder_indices_and_dimensions(slide_content, layout_specs, previous_slide_layout):
    return f"""

        You are an expert AI presentation designer.

Your task is to map slide content to the best PowerPoint slide layout from a given set of layout templates.

---

## Hard Constraints (Must be obeyed):

1. Only choose layouts that contains **all required placeholder types** for the content:
   - A `TITLE` or `CENTER_TITLE` placeholder for the slide title.
   - A `BODY` or `OBJECT` placeholder for bullets. **NEVER use SUBTITLE for bullets.**
   - A `PICTURE` or `OBJECT`placeholder for images, if any.
   - If more than one image if present, choose the one that best fits the layout.
   - Other content (e.g. speaker_notes, captions) must also be mapped to compatible placeholder types.

2. Rank all valid layouts by fitness for this slide based on:
- Placeholder size match for bullets
- Aspect ratio match for images
Pick the top-ranked layout.

3. Do not guess or invent layout names or IDs. Only use layouts from the provided list.

---

## Preferred Layout (Soft Constraints):

1. Use placeholder **size** to improve layout choice:
   - For bullets, prefer placeholders with width ≥ 4.0 inches and height ≥ 2.0 inches.
   - For images, prefer placeholders whose aspect ratio matches the image’s.

2. Use layout aesthetics only **after** the hard rules are satisfied.

---

## Task:

1. **First**, validate which layouts meet all hard constraints.
2. **Second**, select the most appropriate layout from the valid ones based on placeholder sizes and image aspect ratio.
3. **Third**, generate the mapping of content to placeholders in that layout.

---

## Output Format:

Return only a single valid JSON object in this format:

```json
{{
  "slide_number": <int>,
  "layout_id": <int>,
  "layout_name": "<layout_name>",
  "mapping": [
    {{
      "content_type": "title" | "bullets" | "image_path" | "speaker_notes" | "caption",
      "value": "...",
      "placeholder_type": "TITLE" | "BODY" | "OBJECT" | "PICTURE" | ...,
      "placeholder_index": <int>
    }},
    ...
  ]
}}

        Here is the content for the slide:
        {json.dumps(slide_content, indent=2)}

        Here are the available slide layouts:
        {json.dumps(layout_specs, indent=2)}

        previous_slide_layout: {previous_slide_layout}

    """

