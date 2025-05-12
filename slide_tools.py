# slide_tools.py

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

# Create or load a presentation
PPTX_PATH = "output_presentation.pptx"
if os.path.exists(PPTX_PATH):
    prs = Presentation(PPTX_PATH)
else:
    prs = Presentation()

SLIDE_WIDTH = prs.slide_width.inches
SLIDE_HEIGHT = prs.slide_height.inches

def save_presentation():
    prs.save(PPTX_PATH)

def create_slide(title_dict, text_dict, text_box_dict, image_path_dict, image_box_dict):
    slide_layout = prs.slide_layouts[5]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    title = title_dict.get("title")
    text_items = text_dict.get("text", [])
    is_bulleted = text_dict.get("bulleted", False)
    font_size = text_dict.get("font_size", 24)

    text_box = text_box_dict.get("text_box", {})
    image_path = image_path_dict.get("image_path")
    image_box = image_box_dict.get("image_box", {})

    padding = 0.2  # default padding
    text_padding = text_box.get("padding", padding)
    image_padding = image_box.get("padding", padding)
    title_padding = 0.2

    if title:
        title_shape = slide.shapes.title or slide.shapes.add_textbox(Inches(0.5), Inches(title_padding), Inches(9), Inches(0.8))
        title_shape.text = title

    if text_items:
        x = Inches(text_box.get("x") + text_padding)
        y = Inches(text_box.get("y") + text_padding)
        w = Inches(text_box.get("width") - 2 * text_padding)
        h = Inches(text_box.get("height") - 2 * text_padding)
        text_shape = slide.shapes.add_textbox(x, y, w, h)
        tf = text_shape.text_frame
        tf.clear()
        for i, item in enumerate(text_items):
            p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
            p.text = item
            if is_bulleted:
                p.level = 0
                p.text = f"• {p.text}"
            p.font.size = Pt(font_size)
        tf.word_wrap = True

    if image_path and os.path.exists(image_path):
        x = Inches(image_box.get("x") + image_padding)
        y = Inches(image_box.get("y") + image_padding)
        w = Inches(image_box.get("width") - 2 * image_padding)
        h = Inches(image_box.get("height") - 2 * image_padding)
        slide.shapes.add_picture(image_path, x, y, width=w, height=h)

    save_presentation()
    return {
        "status": "success",
        "message": "Slide created with layout positioning and formatting.",
        "title": title,
        "slide_number": len(prs.slides)
    }


def add_text_to_slide(slide_number, text):
    slide = prs.slides[slide_number - 1]
    text_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(5), Inches(2))
    tf = text_box.text_frame
    tf.text = text

    save_presentation()
    return {"status": "success", "message": f"Text added to slide {slide_number}."}

def add_image_to_slide(slide_number, image_path):
    slide = prs.slides[slide_number - 1]
    if os.path.exists(image_path):
        slide.shapes.add_picture(image_path, Inches(5), Inches(1.5), width=Inches(4))
        save_presentation()
        return {"status": "success", "message": f"Image added to slide {slide_number}."}
    return {"status": "error", "message": "Image path not found."}

def add_chart_to_slide(slide_number, chart_data_path):
    # Stub - requires COM or external chart rendering for real implementation
    # For now, this just puts a placeholder box
    slide = prs.slides[slide_number - 1]
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(2), Inches(3), Inches(5), Inches(2)
    )
    shape.text = f"[Chart Placeholder from {chart_data_path}]"
    shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    save_presentation()
    return {"status": "success", "message": f"Chart placeholder added to slide {slide_number}."}


# create_slide(
#     {"title": "Yue Sai's Relevance Decline"},
#     {"text": [
#         "Multiple unsuccessful attempts to reposition the brand",
#         "Failing to adapt to changing consumer preferences",
#         "Losing market share to competitors with more innovative strategies"
#     ], "bulleted": True, "font_size": 20},
#     {"text_box": {
#         "layout": "left", "x": 0.25, "y": 1.0, "width": 6.0, "height": 4.0, "padding": 1
#     }},
#     {"image_path": "images/img_p0_2.png"},
#     {"image_box": {
#         "layout": "right", "x": 5.0, "y": 1.0, "width": 5.0, "height": 4.0, "padding": 1
#     }}
# )

slide_specs = [
    {
        "title": "Targeting Modern, Health-Conscious Young Women in China",
        "text_dict": {
            "text": ['Developing products tailored to health and wellness trends',   
                     'Utilizing digital marketing channels to reach the target demographic',   
                     'Emphasizing natural ingredients and sustainability in product offerings'],
            "bulleted": True,
            "font_size": 16
        },
        "text_box_dict": {
            "text_box": {
                "layout": "right", "x": 0.5, "y": 1.0, "width": 6.0, "height": 6.0, "padding": 0.7
            }
        },
        "image_path_dict": {
            "image_path": "images/img_p0_1.png"
        },
        "image_box_dict": {
            "image_box": {
                "layout": "left", "x": 5, "y": 1.0, "width": 5.0, "height": 4.0, "padding": 0.7
            }
        }
    },
    {
        "title": "Yue Sai's Relevance Decline",
        "text_dict": {
            "text": ["Multiple unsuccessful attempts to reposition the brand",
                        "Failing to adapt to changing consumer preferences",
                        "Losing market share to competitors with more innovative strategies"],
            "bulleted": True,
            "font_size": 18
        },
        "text_box_dict": {
            "text_box": {
                "layout": "left", "x": 0.5, "y": 1.0, "width": 6.0, "height": 6.0, "padding": 0.7
            }
        },
        "image_path_dict": {
            "image_path": "images/img_p0_2.png"
        },
        "image_box_dict": {
            "image_box": {
                "layout": "right", "x": 5.0, "y": 1.0, "width": 5.0, "height": 4.0, "padding": 0.7
            }
        }
    }
]

for spec in slide_specs:
    result = create_slide(
        {"title": spec["title"]},
        spec["text_dict"],
        spec["text_box_dict"],
        spec["image_path_dict"],
        spec["image_box_dict"]
    )
    print(f"✅ Slide created: {result}")